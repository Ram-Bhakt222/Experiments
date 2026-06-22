"""
pipelines/docs/pptx_extract.py — PPTX -> Markdown outline + speaker notes.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict

from pipelines._common import first_input, job_output_dir
from queue_db import update_job


def run(job: Dict[str, Any]) -> None:
    from pptx import Presentation  # python-pptx

    src = first_input(job)
    out = job_output_dir(job)
    update_job(job["id"], progress="reading pptx")

    prs = Presentation(str(src))
    md = []
    md.append(f"# {src.stem}\n")
    for i, slide in enumerate(prs.slides, start=1):
        md.append(f"## Slide {i}")
        # Title shape if any
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text_frame.text or "").strip()
        if title:
            md.append(f"**{title}**\n")
        # Other text shapes
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    for line in t.splitlines():
                        line = line.strip()
                        if line:
                            md.append(f"- {line}")
        # Speaker notes
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    md.append("\n> _Notes:_ " + notes.replace("\n", " ") + "\n")
        except Exception:
            pass
        md.append("")

    md_path = out / f"{src.stem}.md"
    md_path.write_text("\n".join(md), encoding="utf-8")

    update_job(
        job["id"],
        output_path=str(md_path),
        progress=f"done ({len(prs.slides)} slides)",
    )
